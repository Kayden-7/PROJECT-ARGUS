# -*- coding: utf-8 -*-
"""Keep argusdeck(1).pptx's exact design (shapes, fonts, colors, positions —
untouched) and replace its text content with the content from the navy/teal
spec the user pasted (Maya / solo-founder narrative, 10-slide breakdown).
Mapped slide-for-slide onto this deck's existing shape inventory; where a
shape has no direct spec equivalent (e.g. footer tags, level-ladder labels
reused as profile-bullet labels), the connecting text stays in the same
voice as the surrounding spec content rather than inventing new claims.
"""
from pptx import Presentation

SRC = r"C:\Users\baldwin\Downloads\argusdeck(1).pptx"
OUT = r"C:\Users\baldwin\Downloads\argusdeck_v2.pptx"

prs = Presentation(SRC)


def slide(i):
    return prs.slides[i - 1]


def shape(s, sid):
    for sh in s.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_para_runs(s, sid, para_idx, texts):
    tf = shape(s, sid).text_frame
    p = tf.paragraphs[para_idx]
    for r, t in zip(p.runs, texts):
        r.text = t


def set_text(s, sid, text):
    set_para_runs(s, sid, 0, [text])


# ════════════════════════════════ SLIDE 1 ════════════════════════════════════
s = slide(1)
set_text(s, 18, "A Deterministic Permission & Trust Layer for Autonomous AI Agents")
set_text(s, 20, "LLMs Propose. Code Decides.")
set_text(s, 21,
    "We help solo founders who depend on email, and who fear an AI sending or deleting the wrong "
    "thing, finally delegate their inbox to an AI assistant — by solving the fact that AI agents "
    "make their own permission decisions, through a deterministic layer where the AI proposes but "
    "code decides, and trust is earned, so they can let AI handle email without risking a costly mistake.")
set_text(s, 23, "The First Spark Challenge — Challenge A1: Trust & Permissions for Autonomous AI Agents · 25 June 2026")

# ════════════════════════════════ SLIDE 2 ════════════════════════════════════
s = slide(2)
set_text(s, 36, "Challenge A1: Trust & Permissions for Autonomous AI Agents")
set_text(s, 37, "Responding to Challenge A1: Trust & Permissions for Autonomous AI Agents.")
set_para_runs(s, 38, 0, ["AI systems are powerful but uncontrollable — they decide their own scope."])
set_para_runs(s, 38, 2, ["Users need autonomy but also safety — they want help, not surprise."])
set_para_runs(s, 38, 4, ["Today's solutions are all-or-nothing: supervise everything, or hand over full control."])
set_text(s, 40, "AI is powerful but unpredictable. Users want help but need guarantees.")
set_text(s, 41, "How do you give an AI agent the authority to act on your behalf without risking catastrophic mistakes?")
set_para_runs(s, 44, 0, [
    "The First Spark Challenge frames this as ",
    "the core tension",
    " for autonomous agents — and it's exactly the tension ARGUS resolves.",
])

# ════════════════════════════════ SLIDE 3 ════════════════════════════════════
s = slide(3)
set_text(s, 54, "Your User: Solo Founder / Freelancer Who Lives on Email")
set_text(s, 55, "Your user: a solo founder or freelance consultant whose entire livelihood depends on email.")
set_text(s, 57, "ROLE")
set_para_runs(s, 58, 0, ["Role: ", "Solo founder or freelance consultant."])
set_text(s, 74, "CONSTRAINT")
set_para_runs(s, 60, 0, ["Constraint: ", "Entire livelihood depends on email."])
set_text(s, 75, "DAILY REALITY")
set_para_runs(s, 62, 0, ["Daily Reality: ", "80+ emails a day, 4+ hours managing the inbox."])
set_text(s, 76, "THE COST")
set_para_runs(s, 64, 0, ["One mistake costs them: ", "a client relationship, a contract, or revenue."])
set_text(s, 66, "THEIR TYPICAL DAY")
set_text(s, 67,
    "9 AM: a client email demands a same-day reply. 10 AM: another client question needs handling "
    "immediately. 11 AM: an invoice or admin email needs filing correctly. Repeat 10+ times — 4+ "
    "hours of manual work, every day.")
set_text(s, 69, "CORE TENSION")
set_text(s, 70,
    "They want an AI to handle this. But one wrong auto-sent email to the wrong client, or one "
    "deleted thread with contract terms, could cost them real money. So they do it all manually. Stuck.")
set_text(s, 72, "THE STAKES")
set_text(s, 73,
    "They're trapped doing it all manually — drowning in their inbox, but too scared to delegate to "
    "an AI that might make one irreversible mistake. This is exactly who ARGUS is built for.")

# ════════════════════════════════ SLIDE 4 ════════════════════════════════════
s = slide(4)
set_para_runs(s, 88, 0, ["Pain Point", "— what Maya actually fears"])
set_text(s, 91,
    "“I'm afraid an AI will send the wrong thing to the wrong person. I'm terrified it will "
    "delete something important and irreversible. I can't afford to hand over full control — so I "
    "still do it all myself.”")
set_para_runs(s, 92, 0, [
    "The pain is NOT that the AI writes badly — modern tools write well. The pain is ",
    "losing control of an irreversible action",
    ", with no way to undo a mistake before it costs her a client.",
])
set_text(s, 96, "Babysitting the AI — reviewing every single action before it executes — which means no time savings at all.")
set_text(s, 98, "Refusing to use AI altogether — staying overwhelmed, four hours a day, with no help and no scale.")
set_text(s, 100, "One wrong auto-sent email to the wrong client, or one deleted thread with contract terms — gone, irreversible.")
set_text(s, 102, "Either choice keeps her stuck: inefficient but safe, or fast but one mistake from disaster.")
set_para_runs(s, 106, 0, [
    "Maya's actual choice: ",
    "stay inefficient but safe, or trust an AI and risk a catastrophic mistake. There's no third option — until ARGUS.",
])

# ════════════════════════════════ SLIDE 5 (dark panel slide) ════════════════
s = slide(5)
set_para_runs(s, 120, 0, ["The Problem", ": Pain vs. the Real Bottleneck"])
set_text(s, 121, "Pain is what users feel. The bottleneck is the root cause — and the solution must attack the bottleneck, not just the symptom.")
set_text(s, 123, "PAIN (WHAT USERS FEEL)")
set_text(s, 124, "In their own words:")
set_para_runs(s, 125, 0, [
    "They're afraid an AI will send the wrong thing, ",
    "delete something irreversible, or take away their control entirely.",
])
set_text(s, 127, "BOTTLENECK (THE REAL PROBLEM)")
set_text(s, 128,
    "Today's AI agents make their own permission decisions. There's no deterministic layer deciding "
    "what they're allowed to do — it's all-or-nothing: supervise everything, or hand over full "
    "control. There is no middle ground.")
set_text(s, 129, "WHY THIS IS THE ROOT CAUSE, NOT JUST A SYMPTOM")
set_text(s, 131, "Babysit every action and you've built a supervised toy, not an autonomous agent — no time saved at all.")
set_text(s, 133, "Hand over full control instead, and one wrong auto-send or deletion can cost a client relationship or a contract.")
set_text(s, 135, "Refuse to use AI at all, and you're stuck doing everything by hand — no help, no scale, no growth.")
set_para_runs(s, 138, 0, [
    "ARGUS solves the bottleneck: ",
    "by separating intelligence from authority, so users get real delegation without the all-or-nothing risk.",
])

# ════════════════════════════════ SLIDE 6 ════════════════════════════════════
s = slide(6)
set_para_runs(s, 150, 0, ["What Users Do Today", " — and why none of it works"])
set_text(s, 151, "Today, anyone who wants to delegate their inbox to an AI is forced into one of four imperfect options:")
set_text(s, 153, "1. Babysit the AI — review every single action before it executes.")
set_para_runs(s, 154, 0, ["Result: ", "No time savings. She's still doing 90% of the work — defeats the point of delegating."])
set_text(s, 156, "2. Trust the model blindly — let it auto-send and auto-delete without checks.")
set_para_runs(s, 157, 0, ["Result: ", "One wrong send to the wrong client, or one deleted thread, and it's cost her a relationship or a contract."])
set_text(s, 159, "3. Don't use AI at all — keep doing email manually, four hours a day.")
set_para_runs(s, 160, 0, ["Result: ", "Stays overwhelmed, stays trapped. No help, no scale, no growth."])
set_text(s, 162, "4. Hard-coded rules and RPA — traditional deterministic automation.")
set_para_runs(s, 163, 0, ["Result: ", "Deterministic, but brittle and blind — it can't interpret “reply that I'll be there at 1:30” or adapt to context."])
set_para_runs(s, 166, 0, [
    "Maya's actual choice: ",
    "stay inefficient but safe, or trust an AI and risk a catastrophic mistake. There's no third option — until ARGUS.",
])

# ════════════════════════════════ SLIDE 7 ════════════════════════════════════
s = slide(7)
set_text(s, 178, "The Solution: LLMs Propose. Code Decides.")
set_text(s, 179, "ARGUS is a deterministic permission and trust layer between any AI agent and the real actions it wants to take — the layer that lets LLMs propose while code decides.")
set_text(s, 182, "Our core principle is three words:")
set_text(s, 183, "LLMs Propose. Code Decides.")
set_text(s, 184,
    "Layer 1, GPT-4o, interprets your command and produces a structured proposal — recipient, "
    "subject, body, action type. Layer 2, a deterministic Python policy engine, decides if it's "
    "safe: checks trust, validates that approvals are fresh, blocks protected addresses, and outputs "
    "ALLOW, GATED, or BLOCK. Layer 3 is you: see the proposal, approve or reject, and get a window "
    "to undo before it's final.")
set_para_runs(s, 186, 0, [
    "Value proposition, in one line: ",
    "users get the help they need without the all-or-nothing risk — AI proposes, code decides, "
    "trust is earned gradually, and when unsure, it stops and asks.",
])
set_text(s, 188, "ATOMIC")
set_text(s, 189, "All-or-nothing transactions. No partial states — a crash mid-send never leaves a half-sent email or a duplicate.")
set_text(s, 190, "REVERSIBLE")
set_text(s, 191, "A short undo window on every gated action. Change your mind before it's truly final.")
set_text(s, 192, "AUDITED")
set_text(s, 193, "Every decision logged with cryptographic proof — a tamper-evident record of what happened and why.")
set_text(s, 194, "FAILS CLOSED BY CONSTRUCTION")
set_text(s, 195, "Any uncertainty routes to manual review instead of a guess. Safe is the default path, not the lucky one.")
set_text(s, 196, "EARNED TRUST, NOT DECLARED TRUST")
set_text(s, 197, "Autonomy grows with proven reliability and collapses fast after a failure. No “just trust me.”")
set_text(s, 198, "MODEL-AGNOSTIC INFRASTRUCTURE")
set_text(s, 199, "Swap the LLM for any model and the deterministic control plane underneath is unchanged.")

# ════════════════════════════════ SLIDE 8 ════════════════════════════════════
s = slide(8)
set_text(s, 211, "The Architecture: Why It's Fail-Safe")
set_text(s, 212, "Three layers, strictly separated — every decision made by code, never the model.")
set_text(s, 215, "LAYER 1")
set_text(s, 216, "GPT-4o interprets — reads intent, generates a proposal, never decides")
set_text(s, 217,
    "Reads your command — “reply to this client saying I'll send the proposal by Friday” — "
    "and generates a structured proposal: recipient, subject, body, action type. It grounds itself "
    "first, verifying the email it refers to actually exists, and asks for clarification instead of "
    "guessing when something's ambiguous.")
set_text(s, 220, "LAYER 2")
set_text(s, 221, "deterministic Python policy engine — the ONLY thing that can grant permission")
set_text(s, 222,
    "Checks: is this recipient trusted? Is this action ALLOW or GATED? Validates that prior approvals "
    "are still fresh. Blocks protected addresses. Result: “This is GATED — you must approve.” "
    "Not a guess. Not a probability. A decision.")
set_text(s, 225, "LAYER 3")
set_text(s, 226, "you decide — see the proposal, approve or reject, undo window before it's final")
set_text(s, 227,
    "If ALLOW: executes immediately, with audit. If GATED: queued for your approval, with a live "
    "countdown and a window to undo if you change your mind. Every decision is logged with "
    "cryptographic proof — atomic, reversible, audited.")

# ════════════════════════════════ SLIDE 9 ════════════════════════════════════
s = slide(9)
set_para_runs(s, 239, 0, ["What Users Actually Want", " — Survey: N=17"])
set_text(s, 240, "This is not a concept deck — and it's not a guess about what users want, either.")
set_text(s, 242, "88%")
set_text(s, 243, "PREFER PERMISSION-BASED CONTROL")
set_text(s, 245, "100%")
set_text(s, 246, "WANT AI TO PAUSE WHEN UNSURE")
set_text(s, 248, "100%")
set_text(s, 249, "WANT TO SEE WHY, EVERY TIME")
set_text(s, 250, "BACKEND COMPLETE AND WORKING — NINE DELIVERED CAPABILITIES")
set_text(s, 251,
    "Deterministic policy engine · approval queue · earned-trust ledger · crash-safe Gmail execution "
    "· message-style templates · safety filter · append-only hash-chained audit trail · GPT-4o agent "
    "layer · demo mode — plus Phase 8 fail-safes (global emergency stop, and atomic admission with "
    "de-duplication and rate-limiting).")
set_text(s, 252, "869 AUTOMATED TESTS, 100% PASSING — ADVERSARIAL & CHAOS SUITES")
set_text(s, 253,
    "Simulated prompt-injection inside the model's output, mid-send crash recovery, "
    "duplicate-submission storms, and rate-limit abuse. We deliberately attack our own system to "
    "prove it fails closed — safety is tested, not assumed.")
set_text(s, 254, "A LIVE, END-TO-END DEMO — CONNECTED TO A REAL GMAIL ACCOUNT")
set_text(s, 255,
    "Type a command → watch the AI propose → watch deterministic code decide → approve as a human → "
    "see a crash-safe send, with a live audit trail and a trust gauge updating in real time.")
set_text(s, 256, "A FULL WORKING INTERFACE")
set_text(s, 257, "Workbench, a consolidated executions and approval-queue page, an audit trail with a one-click tamper-check, trust history, and settings.")
set_para_runs(s, 260, 0, [
    "Status: ",
    "users didn't ask for a ‘nice AI email tool.’ They asked for exactly what ARGUS builds — backend-complete, fully tested, and demo-ready.",
])

# ════════════════════════════════ SLIDE 10 ═══════════════════════════════════
s = slide(10)
set_para_runs(s, 272, 0, ["We Shipped Real Code.", " Here's the Ask."])
set_text(s, 273, "BY THE NUMBERS")
set_text(s, 274,
    "869 of 869 automated tests passing — 100% coverage. 7 fail-safe controls: hard-stop, rate "
    "limit, dedup, epoch, private contacts, audit chain, stale detection. 3 angles of testing: "
    "normal flow, hacker/adversarial, and strict validation. 0 silent failures — the cryptographic "
    "audit guarantees it. We didn't build a prototype.")
set_text(s, 275, "READY TO SHIP")
set_text(s, 277, "Policy engine locked")
set_text(s, 278, "Phase 8 complete — the deterministic core that decides is done.")
set_text(s, 280, "Live demo running")
set_text(s, 281, "Type a command, watch it propose, watch code decide, approve, watch it send.")
set_text(s, 283, "Full test suite passing")
set_text(s, 284, "869 of 869, including the adversarial and chaos suites.")
set_text(s, 286, "Architecture scales beyond email")
set_text(s, 287, "Calendars, CRMs, file systems — anywhere an agent takes an irreversible action.")
set_text(s, 288, "The close")
set_text(s, 289, "AI should amplify your judgment, not replace it. ARGUS is that AI.")
set_text(s, 291, "Live Demo: localhost:8081  ·  GitHub: [repo link]  ·  Contact: kayden.low24@gmail.com")
set_text(s, 292, "LLMs Propose. Code Decides.")

prs.save(OUT)
print("Saved:", OUT)
