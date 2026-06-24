"""Replace slide 10 (currently 'The Pitch in One Line + Judge Check') with a
new 'Vision & The Ask' slide: grounded in ARGUS's real ICP (the solo founder /
freelance consultant who won't hand over their inbox), then the forward vision
(Claude-native rebuild, real product) and the ask. Same shapes, same theme,
same positions — only text content changes, run-for-run, preserving each
run's existing color/font/size so the visual design is untouched.
"""
from pptx import Presentation

SRC = r"C:\Users\baldwin\Downloads\Website_redesign_deck_10slides.pptx"
OUT = r"C:\Users\baldwin\Downloads\Website_redesign_deck_10slides_v2.pptx"

prs = Presentation(SRC)
s = prs.slides[9]


def shape(name):
    for sh in s.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def set_runs(name, texts):
    """Set text run-by-run, reusing each existing run's formatting (color/
    font/size untouched) — texts must match the existing run count."""
    tf = shape(name).text_frame
    p0 = tf.paragraphs[0]
    assert len(p0.runs) == len(texts), f"{name}: {len(p0.runs)} runs vs {len(texts)} texts"
    for r, t in zip(p0.runs, texts):
        r.text = t
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def set_text(name, text):
    set_runs(name, [text])


# ── Title (2 runs: white hook, grey continuation) ──────────────────────────
set_runs("Text 6", [
    "Built for the One Email You Can't Get Wrong.",
    " Now, the Vision Beyond It.",
])

# ── Left column: who we built this for ─────────────────────────────────────
set_text("Text 7", "WHO WE BUILT THIS FOR")
set_runs("Text 8", [
    "A ",
    "solo founder or freelance consultant",
    " whose livelihood runs on email — who wants an AI assistant to handle the inbox, but ",
    "won't hand it over",
    " because one wrong auto-sent message to a client, or one deleted thread, could cost ",
    "a relationship or a contract",
    ". So today, they still do it all by hand — and ARGUS exists so they finally don't have to.",
])

# ── Right column: where ARGUS goes from here ────────────────────────────────
set_text("Text 9", "WHERE ARGUS GOES FROM HERE")

set_text("Text 11", "Today")
set_text("Text 12", "A real, working product, not a slide — live Gmail OAuth, a deterministic policy engine, and 863 passing automated tests.")

set_text("Text 14", "The architecture already plans for this")
set_text("Text 15", "The PROPOSE layer was built model-agnostic — swap the model, the deterministic core stays the same.")

set_text("Text 17", "The rebuild")
set_text("Text 18", "The next build is Claude-native — Claude reasons, the deterministic core still decides every action, unchanged.")

set_text("Text 20", "The ask")
set_text("Text 21", "Help us take ARGUS past the hackathon: harden the Claude-native build for real inboxes.")

set_text("Text 22", "Why now")
set_text("Text 23", "Every agent platform hits the same wall: authority without losing control. We intend to build the layer that solves it.")

# ── Bottom band: close ──────────────────────────────────────────────────────
set_text("Text 25", "This began as a hackathon build. The vision: a real product, Claude-native, trusted by the people who need it.")
# Text 26 ("AI proposes. Code decides.") stays unchanged — the deck's recurring tagline.

prs.save(OUT)
print("Saved:", OUT)
