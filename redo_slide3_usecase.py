# -*- coding: utf-8 -*-
"""Slide 3 (Target User & Use Case): keep the Level 1-4 ladder (left column)
completely untouched. Replace the right column's three boxes — currently
"THE SPECIFIC MOMENT" / "FIRST TARGET USER" / "EXPANSION MARKET" — with new,
standalone content about the use case and the target user, written fresh
(not reused from any other slide in the deck). Design/position/font/color
untouched, text only.
"""
import math
from pptx import Presentation
from pptx.util import Pt

SRC = r"C:\Users\baldwin\Downloads\argusdeck(1).pptx"
OUT = r"C:\Users\baldwin\Downloads\argusdeck(1)_slide3_v2.pptx"

prs = Presentation(SRC)
s = prs.slides[2]
MIN_SCALE = 0.62


def shape(sid):
    for sh in s.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_text(sid, new_text, min_scale=MIN_SCALE):
    sh = shape(sid)
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    r0 = p0.runs[0]
    orig_len = len(tf.text) or 1
    orig_size = r0.font.size.pt if r0.font.size else 18.0
    new_len = len(new_text)
    scale = 1.0
    if new_len > orig_len:
        scale = max(min_scale, math.sqrt(orig_len / new_len))
    r0.text = new_text
    r0.font.size = Pt(round(orig_size * scale, 1))
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


# Left column (Level 1-4 ladder): ids 57,58,60,62,64,74,75,76 — untouched.

# Right column — replaced with fresh use-case / target-user content.
set_text(66, "THE USE CASE")
set_text(67,
    "A client emails asking to move up a deadline. The founder tells ARGUS what to reply. ARGUS "
    "drafts it, checks whether an external send needs a human yes, and either fires it off or "
    "waits for one tap of approval — the same flow for every email, every day, without the "
    "founder reading every draft first.")

set_text(69, "WHAT THEY'RE WILLING TO TRY")
set_text(70,
    "Not full autonomy on day one — a system that earns the right to act on more, one proven "
    "action at a time.")

set_text(72, "WHO THEY ARE")
set_text(73,
    "There is no assistant, no ops team — they are the founder, the closer, and the last line "
    "of defence on every send, alone.")

prs.save(OUT)
print("Saved:", OUT)
