# -*- coding: utf-8 -*-
"""Replace ALL content in argusdeck(1).pptx with the new Maya/solo-founder
spec, while keeping every shape, position, color, and font as-is. Existing
slides 4+5 (Pain, Bottleneck) absorb the new merged "Pain vs Bottleneck"
topic between them; existing slide 5's leftover capacity is repurposed for
the new "Live Demo" topic so all 10 new topics still land across 10 slides
without resizing or deleting any shape. Font size is auto-scaled down only
when new text is longer than the original (never scaled up), so headings
stay bigger than body text and nothing overflows its box.
"""
import math
from pptx import Presentation
from pptx.util import Pt

SRC = r"C:\Users\baldwin\Downloads\argusdeck(1).pptx"
OUT = r"C:\Users\baldwin\Downloads\argusdeck(1)_v2.pptx"

prs = Presentation(SRC)
MIN_SCALE = 0.62


def slide(i):
    return prs.slides[i - 1]


def shape(sl, sid):
    for sh in sl.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_text(sl, sid, new_text, min_scale=MIN_SCALE):
    sh = shape(sl, sid)
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    r0 = p0.runs[0]
    orig_len = len(tf.text) or 1
    orig_size = r0.font.size.pt if r0.font.size else 18.0
    new_len = len(new_text)
    scale = 1.0
    if new_len > orig_len:
        scale = max(min_scale, math.sqrt(orig_len / new_len))
    new_size = round(orig_size * scale, 1)
    r0.text = new_text
    r0.font.size = Pt(new_size)
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def set_multiline(sl, sid, lines, min_scale=MIN_SCALE):
    """Same as set_text but writes each line as its own paragraph, cloning
    the first paragraph's run formatting onto every new paragraph."""
    sh = shape(sl, sid)
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    r0 = p0.runs[0]
    orig_len = len(tf.text) or 1
    orig_size = r0.font.size.pt if r0.font.size else 18.0
    joined = " ".join(lines)
    new_len = len(joined)
    scale = 1.0
    if new_len > orig_len:
        scale = max(min_scale, math.sqrt(orig_len / new_len))
    new_size = round(orig_size * scale, 1)
    font_name, bold, italic, color = r0.font.name, r0.font.bold, r0.font.italic, None
    try:
        color = r0.font.color.rgb
    except Exception:
        pass
    r0.text = lines[0]
    r0.font.size = Pt(new_size)
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(new_size)
        r.font.name = font_name
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color


# ════════════════════════════════ SLIDE 1 — TITLE ═══════════════════════════
s = slide(1)
set_text(s, 18, "A Deterministic Permission & Trust Layer for Autonomous AI Agents")
set_text(s, 20, "LLMs Propose. Code Decides.")
set_text(s, 21,
    "We help solo founders who depend on email who struggle with fear of an AI sending or "
    "deleting the wrong thing during delegating their inbox to an AI assistant by solving the "
    "fact that AI agents make their own permission decisions through a deterministic layer where "
    "the AI proposes but code decides and trust is earned so they can finally let AI handle email "
    "without risking a costly mistake.")
set_text(s, 23, "The First Spark Challenge — Challenge A1: Trust & Permissions for Autonomous AI Agents · 25 June 2026")

# ════════════════════════════════ SLIDE 2 — THE CHALLENGE ═══════════════════
s = slide(2)
set_text(s, 36, "Challenge A1: Trust & Permissions for Autonomous AI Agents")
set_text(s, 37, "The core tension every autonomous-agent product must resolve.")
set_multiline(s, 38, [
    "AI systems are powerful but uncontrollable because they decide their own scope.",
    "Users need autonomy but also safety because they want help, not surprise.",
    "Today's solutions are all-or-nothing: supervise everything or hand over full control.",
])
set_text(s, 40, "Before any action executes, one question has to be answered:")
set_text(s, 41, "How do you give an AI agent the authority to act on your behalf without risking catastrophic mistakes?")
set_text(s, 44, "The First Spark Challenge frames this as the core tension for autonomous agents.")

# ════════════════════════════════ SLIDE 3 — TARGET USER ═════════════════════
s = slide(3)
set_text(s, 54, "Your User: Solo Founder / Freelancer Who Lives on Email")
set_text(s, 55, "A solo founder or freelance consultant who lives in their inbox — one bad send away from losing a client.")
set_text(s, 57, "ROLE")
set_text(s, 58, "Solo founder or freelance consultant.")
set_text(s, 74, "CONSTRAINT")
set_text(s, 60, "Their entire livelihood depends on email.")
set_text(s, 75, "DAILY REALITY")
set_text(s, 62, "80+ emails a day, 4+ hours managing the inbox.")
set_text(s, 76, "CORE TENSION")
set_text(s, 64,
    "They want an AI to handle this. But one wrong auto-sent email to the wrong client, or one "
    "deleted thread with contract terms, could cost them a relationship, a contract, or real "
    "revenue. So they do it all manually. Stuck.")
set_text(s, 66, "9 AM")
set_text(s, 67, "Client email arrives asking for a proposal by end of day — must reply fast.")
set_text(s, 69, "10 AM")
set_text(s, 70, "Another client question arrives — must handle it immediately.")
set_text(s, 72, "11 AM, THEN REPEAT")
set_text(s, 73,
    "Invoice or admin email needs filing — and this repeats 10+ times, costing 4+ hours of "
    "manual work a day.")

# ═══════════════════ SLIDE 4 — PAIN + BOTTLENECK (merged, split across this slide's own shapes) ═══
s = slide(4)
set_text(s, 88, "The Problem: Two Layers")
set_text(s, 89, "PAIN: WHAT USERS FEEL")
set_text(s, 91,
    "“I'm afraid an AI will send the wrong thing to the wrong person.” “I'm terrified it will "
    "delete something important and irreversible.” “I can't afford to hand over full control.”")
set_text(s, 92,
    "The bottleneck: today's AI agents make their own permission decisions — there is no "
    "deterministic layer deciding what they are allowed to do.")
set_text(s, 93, "BOTTLENECK: THE REAL PROBLEM")
set_text(s, 96, "It is all-or-nothing: either users supervise everything or the AI has full autonomy.")
set_text(s, 98, "There is no deterministic layer deciding what they're allowed to do.")
set_text(s, 100, "Today's AI agents make their own permission decisions.")
set_text(s, 102, "There is no middle ground.")
set_text(s, 106, "ARGUS solves the bottleneck by separating intelligence from authority.")

# ════════════════════════════════ SLIDE 5 — repurposed as LIVE DEMO ══════════
s = slide(5)
set_text(s, 120, "See It In Action")
set_text(s, 121, "A live, end-to-end walk-through — five steps from inbox to a crash-safe send.")
set_text(s, 123, "STEP 1 — INBOX")
set_text(s, 124, "Select an email from a colleague.")
set_text(s, 125, "Step 2 — Command: type your intent, e.g. “reply that I'll send the proposal by Friday.”")
set_text(s, 127, "STEP 3 — PROPOSAL")
set_text(s, 128,
    "ARGUS generates a draft. Step 4 — Queue: approve it, with a 4:00 undo countdown before it "
    "sends. Step 5 — Executed: a toast confirms the send.")
set_text(s, 129, "LIVE DEMO ACCESS")
set_text(s, 131, "Live Demo URL: localhost:8081")
set_text(s, 133, "Credentials: PROJECT_ARGUS / ARGUS_DEMO")
set_text(s, 135, "Features available: emergency stop, profile switching, audit trail inspection.")
set_text(s, 138, "Live demo is running during the pitch. If connection fails, these steps are the fallback.")

# ════════════════════════════════ SLIDE 6 — CURRENT ALTERNATIVE ═════════════
s = slide(6)
set_text(s, 150, "What Users Do Today (And Why It Doesn't Work)")
set_text(s, 151, "Today, anyone who wants to delegate to an AI agent is stuck choosing between two broken paths.")
set_text(s, 153, "Path A: Babysit the AI")
set_text(s, 154, "Review every action before it executes — no real time savings, still doing 90% of the work.")
set_text(s, 156, "Path B: Don't Use AI At All")
set_text(s, 157, "Keep doing email manually, 4 hours a day — stays overwhelmed, stays trapped. No help, no scale, no growth.")
set_text(s, 159, "Maya's Actual Choice")
set_text(s, 160, "Stay inefficient but safe, or trust an AI and risk a catastrophic mistake.")
set_text(s, 162, "The Missing Option")
set_text(s, 163, "There is no third option — until ARGUS.")
set_text(s, 166,
    "Both paths cost her the same thing — her time, or her safety. ARGUS is the third option "
    "neither one offers.")

# ════════════════════════════════ SLIDE 7 — SOLUTION & VALUE ═══════════════
s = slide(7)
set_text(s, 178, "The Solution: LLMs Propose. Code Decides.")
set_text(s, 179,
    "ARGUS is the deterministic layer that sits between the AI and the inbox — every action "
    "checked before it ever executes.")
set_text(s, 182, "Our core principle is three words:")
set_text(s, 183, "LLMs Propose. Code Decides.")
set_text(s, 184,
    "Layer 1, GPT-4o, reads the user's command — “reply to this client saying I'll send the "
    "proposal by Friday” — and generates a structured proposal: recipient, subject, body, action "
    "type. Layer 2, the deterministic Python policy engine, checks whether the recipient is "
    "trusted and whether the action is ALLOW or GATED, validates that all prior approvals are "
    "still fresh via the hard-stop epoch, and confirms the address isn't protected by the "
    "private-contact guard.")
set_text(s, 186,
    "Layer 3: the user sees the proposal, approves or rejects it, gets a 4-minute undo window, "
    "and a full audit trail showing exactly why the code decided GATED. The AI proposes. Code "
    "decides. Trust is earned gradually — and reversible, not just approved.")
set_text(s, 187, "WHY IT WORKS")
set_text(s, 188, "GETS HELP WITHOUT ALL-OR-NOTHING RISK")
set_text(s, 189, "She delegates real email actions without handing over full control.")
set_text(s, 190, "WHEN UNSURE, IT STOPS AND ASKS")
set_text(s, 191, "Any uncertain action pauses for her decision instead of guessing.")
set_text(s, 192, "TRUST IS EARNED, NOT ASSUMED")
set_text(s, 193, "Autonomy grows gradually as the system proves it gets things right.")
set_text(s, 194, "REVERSIBLE, NOT JUST APPROVED")
set_text(s, 195, "A 4-minute undo window means even an approved send can be pulled back.")
set_text(s, 196, "EVERY DECISION IS AUDITABLE")
set_text(s, 197, "She can always see exactly why the code allowed or gated an action.")
set_text(s, 198, "THE AI NEVER GETS THE LAST WORD")
set_text(s, 199, "GPT-4o proposes. Deterministic code is the only thing that decides.")

# ════════════════════════════════ SLIDE 8 — HOW IT WORKS ════════════════════
s = slide(8)
set_text(s, 211, "The Architecture: Why It's Fail-Safe")
set_text(s, 212, "Email becomes a command. Three layers turn that command into a safe, auditable action.")
set_text(s, 215, "LAYER 1")
set_text(s, 216, "GPT-4o · interprets intent")
set_text(s, 217,
    "Reads the user's command — “reply to this client saying I'll send the proposal by Friday” "
    "— and generates a structured proposal: recipient, subject, body, and action type.")
set_text(s, 220, "LAYER 2")
set_text(s, 221, "Deterministic policy engine — the ONLY thing that can grant permission")
set_text(s, 222,
    "Checks whether the recipient is trusted and whether the action is ALLOW or GATED. Validates "
    "that all prior approvals are still fresh via the hard-stop epoch. Confirms the address isn't "
    "blocked by the private-contact guard. Also checks for duplicate submissions and rate-limit "
    "abuse. Decision: ALLOW, GATED, or BLOCK.")
set_text(s, 225, "LAYER 3")
set_text(s, 226, "User decision, then atomic, crash-safe Gmail execution")
set_text(s, 227,
    "If ALLOW, it executes immediately with a full audit trail. If GATED, it queues for her "
    "approval with a 4-minute undo window. Every send is atomic — no partial states — and every "
    "decision is logged with cryptographic, tamper-evident proof.")

# ════════════════════════════════ SLIDE 9 — VALIDATION (SURVEY) ═════════════
s = slide(9)
set_text(s, 239, "What Users Actually Want (Survey: N=17)")
set_text(s, 240, "We surveyed 17 users about handing email over to an AI agent. The pattern was unanimous.")
set_text(s, 242, "88%")
set_text(s, 243, "PREFER “ASK EVERY TIME”")
set_text(s, 245, "100%")
set_text(s, 246, "PAUSE WHEN UNSURE")
set_text(s, 248, "100%")
set_text(s, 249, "WANT TO SEE WHY")
set_text(s, 250, "WHY: “ASK PERMISSION,” NOT “EARN AUTONOMY”")
set_text(s, 251,
    "88% prefer the AI ask permission every time, versus letting it earn autonomy gradually on "
    "its own.")
set_text(s, 252, "WHY: PAUSE, DON'T GUESS")
set_text(s, 253,
    "100% want the AI to pause when it's unsure, rather than make its best guess and act anyway.")
set_text(s, 254, "WHY: TRANSPARENCY, NOT A BLACK BOX")
set_text(s, 255,
    "100% want to see exactly why each decision was made — in their words, the AI should "
    "“justify its every decision.”")
set_text(s, 256, "THE BOTTOM LINE")
set_text(s, 257, "Users didn't ask for a “nice AI email tool.” They asked for exactly what ARGUS builds.")
set_text(s, 260, "We didn't invent demand for ARGUS. We solved for it.")

# ════════════════════════════════ SLIDE 10 — PROOF OF EXECUTION & CLOSE ═════
s = slide(10)
set_text(s, 272, "We Shipped Real Code")
set_text(s, 273, "BY THE NUMBERS")
set_text(s, 274,
    "869 out of 869 tests passing — 100% coverage. Seven fail-safe controls: hard-stop, rate "
    "limit, dedup, epoch, private contacts, audit chain, and stale detection. Three angles of "
    "testing — normal flow, hacker/adversarial, and strict validation. Zero silent failures, "
    "because the cryptographic audit guarantees transparency.")
set_text(s, 275, "READY TO SHIP")
set_text(s, 277, "Locked")
set_text(s, 278, "Policy engine locked — Phase 8 complete.")
set_text(s, 280, "Running")
set_text(s, 281, "Live demo running, connected to a real account.")
set_text(s, 283, "Passing")
set_text(s, 284, "Full automated test suite passing, top to bottom.")
set_text(s, 286, "Scalable")
set_text(s, 287, "Architecture already scales beyond a single inbox.")
set_text(s, 288, "Closing")
set_text(s, 289, "AI should amplify your judgment, not replace it. ARGUS is that AI.")
set_text(s, 291, "Live Demo: localhost:8081 · GitHub: [repo link] · Contact: kayden.low24@gmail.com")
set_text(s, 292, "LLMs Propose. Code Decides.")

prs.save(OUT)
print("Final save complete (all 10 slides):", OUT)
