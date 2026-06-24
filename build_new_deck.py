# -*- coding: utf-8 -*-
"""Build the new 10-slide ARGUS deck from scratch per the navy/teal design
spec (Maya persona, 869/869 tests, N=17 survey). Separate file from dck.pptx.
Icons are emoji glyphs (as the spec allows) rather than external SVGs since
no icon library is wired into this environment. Slide 8's "real screenshots"
are placeholders (captioned text boxes) — no live-demo screenshots were
provided to embed. The GitHub link is left as the literal placeholder
"[repo link]" from the spec since no real URL was given.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r"C:\Users\baldwin\Downloads\ARGUS_NewDeck_NavyTeal.pptx"

W, H = 13.333, 7.5

NAVY = "1A3A52"
TEAL = "00BFA5"
GREEN = "2E7D32"
AMBER = "F57C00"
RED = "D32F2F"
WHITE = "FFFFFF"
LGRAY_TXT = "CCCCCC"
DGRAY = "333333"
GRAY12 = "666666"

BG_TEAL_LT = "E0F7F6"
BG_ORANGE_LT = "FFF3E0"
BG_GREEN_LT = "E8F5E9"
BG_BLUE_LT = "E3F2FD"
BG_RED_LT = "FFEBEE"
BG_GRAY_LT = "FAFAFA"
BG_YELLOWTEAL_LT = "F0F4C3"

BORDER_ORANGE = "FFCC80"
BORDER_GREEN = "81C784"
BORDER_RED = "EF5350"
BORDER_BLUE = "64B5F6"
BLUE = "1976D2"

SERIF = "Georgia"
SANS = "Arial"
MONO = "Courier New"


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


def px2pt(px):
    return Pt(round(px * 0.75, 1))


def pct(p, total):
    return Inches(total * p / 100.0)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def new_slide(bg_hex=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(bg_hex)
    return s


def add_box(slide, x, y, w, h, fill_hex=None, line_hex=None, line_w=1.5, dashed=False, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill_hex:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill_hex)
    else:
        shp.fill.background()
    if line_hex:
        shp.line.color.rgb = rgb(line_hex)
        shp.line.width = Pt(line_w)
        if dashed:
            ln = shp.line._get_or_add_ln()
            dash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
            ln.append(dash)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size_px, color_hex, font=SANS, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing != 1.0:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = px2pt(size_px)
        r.font.name = font
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = rgb(color_hex)
    return tb


def add_bullets(slide, x, y, w, h, items, size_px, color_hex, font=SANS, bold=False,
                 italic=False, bullet=True, line_spacing=1.3, space_after_px=6):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        txt, ibold = (item, bold) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = px2pt(space_after_px)
        r = p.add_run()
        r.text = (u"•  " if bullet else "") + txt
        r.font.size = px2pt(size_px)
        r.font.name = font
        r.font.bold = ibold
        r.font.italic = italic
        r.font.color.rgb = rgb(color_hex)
    return tb


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ════════════════════════════════ SLIDE 1 — TITLE ═══════════════════════════
s = new_slide(NAVY)
add_text(s, 0, H * 0.06, W, H * 0.10, "ARGUS", 24, TEAL, font=MONO, bold=True, align=PP_ALIGN.CENTER)
add_text(s, W * 0.08, H * 0.30, W * 0.84, H * 0.28,
         "A Deterministic Permission & Trust Layer for Autonomous AI Agents",
         44, WHITE, font=SERIF, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
add_box(s, W * 0.30, H * 0.595, W * 0.40, 0.02, fill_hex=TEAL)
add_text(s, W * 0.08, H * 0.62, W * 0.84, H * 0.13,
         "LLMs Propose. Code Decides.", 36, TEAL, font=SERIF, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, W * 0.15, H * 0.77, W * 0.70, H * 0.13,
         "We help solo founders who depend on email and fear an AI sending or deleting the wrong "
         "thing delegate their inbox to an AI assistant — by solving the fact that AI agents make "
         "their own permission decisions, through a deterministic layer where the AI proposes but "
         "code decides, and trust is earned, so they can finally let AI handle email without risking "
         "a costly mistake.",
         11, LGRAY_TXT, font=SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
add_text(s, 0, H * 0.93, W, H * 0.06,
         "The First Spark Challenge — Challenge A1: Trust & Permissions for Autonomous AI Agents · 25 June 2026",
         10, LGRAY_TXT, font=SANS, align=PP_ALIGN.CENTER)
set_notes(s, "This is ARGUS. We exist to solve one specific problem: how do you let an AI agent handle "
             "your email without risking an irreversible mistake? The answer is radical transparency, "
             "deterministic control, and earned trust.")

# ════════════════════════════════ SLIDE 2 — THE CHALLENGE ═══════════════════
s = new_slide(WHITE)
add_text(s, 0.6, H * 0.06, W - 1.2, H * 0.10, "Challenge A1: Trust & Permissions for Autonomous AI Agents",
         32, NAVY, font=SERIF, bold=True, align=PP_ALIGN.LEFT)
add_text(s, 0.8, H * 0.25, W - 1.6, H * 0.22,
         "How do you give an AI agent the authority to act on your behalf without risking "
         "catastrophic mistakes?",
         26, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
add_bullets(s, 0.8, H * 0.55, W * 0.55, H * 0.20, [
    "AI systems are powerful but uncontrollable (they decide their own scope)",
    "Users need autonomy but also safety (they want help, not surprise)",
    "Today's solutions are all-or-nothing: supervise everything OR hand over full control",
], 16, DGRAY, line_spacing=1.3)
add_text(s, W * 0.62, H * 0.30, W * 0.30, H * 0.30, u"⚖", 96, TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0, H * 0.94, W, H * 0.05, "The First Spark Challenge frames this as the core tension for autonomous agents.",
         11, GRAY12, italic=True, align=PP_ALIGN.CENTER)
set_notes(s, "The First Spark Challenge asks: how do you trust an AI to act autonomously? This is the "
             "core tension. AI is powerful but unpredictable. Users want help but need guarantees. "
             "ARGUS is the answer to that exact question.")

# ════════════════════════════════ SLIDE 3 — TARGET USER ═════════════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.06, W, H * 0.09, "Your User: Solo Founder / Freelancer Who Lives on Email",
         30, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0.6, H * 0.22, 4, H * 0.04, "User Profile", 15, NAVY, bold=True)
add_bullets(s, 0.6, H * 0.265, W * 0.55, H * 0.18, [
    "Role: Solo founder or freelance consultant",
    "Constraint: Entire livelihood depends on email",
    "Daily Reality: 80+ emails/day, 4+ hours managing inbox",
    "One mistake costs them: A client relationship, a contract, or revenue",
], 15, DGRAY, line_spacing=1.25, space_after_px=4)
add_text(s, 0.6, H * 0.50, 4, H * 0.04, "Their Typical Day", 15, NAVY, bold=True)
add_bullets(s, 0.6, H * 0.545, W * 0.55, H * 0.18, [
    "9 AM: Client email arrives asking for proposal by EOD → Must reply fast",
    "10 AM: Another client question → Must handle immediately",
    "11 AM: Invoice/admin email → Must file correctly",
    ("Repeat 10+ times → 4+ hours of manual work", True),
], 14, DGRAY, line_spacing=1.2, space_after_px=4)
box = add_box(s, 0.6, H * 0.76, W * 0.55, H * 0.16, fill_hex=BG_TEAL_LT)
box.line.fill.background()
left_bar = add_box(s, 0.6, H * 0.76, 0.06, H * 0.16, fill_hex=TEAL)
add_text(s, 0.8, H * 0.765, W * 0.50, H * 0.15,
         "They want an AI to handle this. But one wrong auto-sent email to the wrong client, or one "
         "deleted thread with contract terms, could cost them real money. So they do it all manually. Stuck.",
         13.5, DGRAY, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
add_text(s, W * 0.68, H * 0.30, W * 0.28, H * 0.30, u"\U0001F4E7", 90, TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "Your user is a solo founder or freelancer whose entire livelihood depends on email. They "
             "get 80+ emails a day. Each one needs a reply, a decision, or an action. They want AI to "
             "help. But they can't afford a single mistake — one auto-sent email to the wrong client "
             "could cost them a month's revenue. So they're trapped doing it all manually. Drowning but "
             "too scared to delegate. This is who ARGUS is for.")

# ════════════════════════════════ SLIDE 4 — PAIN vs BOTTLENECK ══════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.05, W, H * 0.09, "The Problem: Two Layers", 30, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)
add_box(s, 0.6, H * 0.22, W * 0.40, H * 0.55, fill_hex=BG_ORANGE_LT, line_hex=BORDER_ORANGE, line_w=1)
add_text(s, 0.9, H * 0.26, W * 0.34, H * 0.05, "PAIN (What Users Feel)", 16, AMBER, bold=True)
add_bullets(s, 0.9, H * 0.33, W * 0.34, H * 0.40, [
    "“I'm afraid an AI will send the wrong thing to the wrong person”",
    "“I'm terrified it will delete something important and irreversible”",
    "“I can't afford to hand over full control”",
], 15, DGRAY, line_spacing=1.3)
add_text(s, W * 0.42, H * 0.30, W * 0.16, H * 0.06, "ROOT CAUSE", 11, NAVY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, W * 0.42, H * 0.38, W * 0.16, H * 0.10, u"→", 30, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_box(s, W * 0.58, H * 0.22, W * 0.40, H * 0.55, fill_hex=BG_GREEN_LT, line_hex=BORDER_GREEN, line_w=1)
add_text(s, W * 0.61, H * 0.26, W * 0.34, H * 0.05, "BOTTLENECK (The Real Problem)", 16, GREEN, bold=True)
add_bullets(s, W * 0.61, H * 0.33, W * 0.34, H * 0.40, [
    ("Today's AI agents make their own permission decisions", True),
    "No deterministic layer deciding what they're allowed to do",
    "It's all-or-nothing: either users supervise everything OR the AI has full autonomy",
    ("There is no middle ground.", True),
], 15, DGRAY, line_spacing=1.3)
add_text(s, 0, H * 0.83, W, H * 0.06, "ARGUS solves the bottleneck by separating intelligence from authority.",
         13, NAVY, italic=True, align=PP_ALIGN.CENTER)
set_notes(s, "Pain is what users feel: fear. But the BOTTLENECK — the root cause — is that AI agents "
             "decide their own permissions. Today's AI tools are binary. Either users watch every single "
             "action (which wastes their time) or they hand over full control (which is too scary). "
             "There's no middle ground. No deterministic layer that says 'this action is safe to propose' "
             "or 'this is too risky, ask the user first.' ARGUS solves that bottleneck.")

# ════════════════════════════════ SLIDE 5 — CURRENT ALTERNATIVES ════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.05, W, H * 0.09, "What Users Do Today (And Why It Doesn't Work)",
         28, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)
add_box(s, 0.6, H * 0.22, W * 0.42, H * 0.52, fill_hex=BG_RED_LT, line_hex=BORDER_RED, line_w=1.75, dashed=True)
add_text(s, 0.9, H * 0.255, W * 0.36, H * 0.05, "Path A: Babysit the AI", 16, RED, bold=True)
add_text(s, 0.9, H * 0.325, W * 0.36, H * 0.08, "Use an AI email tool but review every action before it executes",
         14, DGRAY, line_spacing=1.25)
add_text(s, 0.9, H * 0.46, W * 0.36, H * 0.07, "Result: No time savings. She's still doing 90% of the work.",
         14, RED, bold=True, line_spacing=1.2)
add_text(s, 0.9, H * 0.58, W * 0.36, H * 0.06, "Defeats the point of delegating.", 13, DGRAY, italic=True)
add_text(s, W * 0.33, H * 0.25, 1.2, 0.7, u"✗", 36, RED, align=PP_ALIGN.RIGHT)
add_box(s, W * 0.55, H * 0.22, W * 0.42, H * 0.52, fill_hex=BG_RED_LT, line_hex=BORDER_RED, line_w=1.75, dashed=True)
add_text(s, W * 0.58, H * 0.255, W * 0.36, H * 0.05, "Path B: Don't Use AI At All", 16, RED, bold=True)
add_text(s, W * 0.58, H * 0.325, W * 0.36, H * 0.08, "Keep doing email manually, 4 hours/day", 14, DGRAY, line_spacing=1.25)
add_text(s, W * 0.58, H * 0.46, W * 0.36, H * 0.07, "Result: Stays overwhelmed, stays trapped.", 14, RED, bold=True)
add_text(s, W * 0.58, H * 0.58, W * 0.36, H * 0.06, "No help, no scale, no growth.", 13, DGRAY, italic=True)
add_text(s, W * 0.89, H * 0.25, 1.2, 0.7, u"✗", 36, RED, align=PP_ALIGN.RIGHT)
add_text(s, 0, H * 0.80, W, H * 0.07,
         "Maya's actual choice: stay inefficient but safe, OR trust an AI and risk a catastrophic mistake.",
         16, TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0, H * 0.89, W, H * 0.06, "There's no third option — until ARGUS.",
         13, NAVY, italic=True, align=PP_ALIGN.CENTER)
set_notes(s, "Right now, the user is stuck. They can babysit an AI (which saves no time) or refuse to "
             "use one (which saves no time either). Both fail. They need a third option: an AI they can "
             "genuinely delegate to without risking a costly mistake. That's the gap ARGUS fills.")

# ════════════════════════════════ SLIDE 6 — SOLUTION & VALUE PROP ═══════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.04, W, H * 0.08, "The Solution: LLMs Propose. Code Decides.",
         28, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)


def layer_box(y_pct, h_pct, bg, border, icon, title, title_color, lines):
    add_box(s, 0.7, H * y_pct, W - 1.4, H * h_pct, fill_hex=bg, line_hex=border, line_w=1.75)
    add_text(s, 0.95, H * y_pct + 0.04, 0.6, H * h_pct - 0.08, icon, 26, title_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.6, H * y_pct + 0.06, W - 2.5, 0.3, title, 15, title_color, bold=True)
    add_bullets(s, 1.6, H * y_pct + 0.36, W - 2.5, H * h_pct - 0.42, lines, 12.5, DGRAY, line_spacing=1.2, space_after_px=2)


layer_box(0.135, 0.155, BG_BLUE_LT, BLUE, u"\U0001F9E0", "Layer 1: GPT-4o Interprets", BLUE, [
    "Reads user's command: “Reply to this client saying I'll send the proposal by Friday”",
    "Generates a structured proposal: recipient, subject, body, action type",
])
add_text(s, W / 2 - 0.3, H * 0.295, 0.6, H * 0.03, u"↓", 16, NAVY, bold=True, align=PP_ALIGN.CENTER)
layer_box(0.33, 0.195, BG_GREEN_LT, GREEN, u"⚙", "Layer 2: Deterministic Code Decides", GREEN, [
    "Python policy engine checks: Is this recipient trusted? Is this action ALLOW or GATED?",
    "Hard-stop epoch validates: Are all prior approvals still fresh?",
    "Private-contact guard: Is this address protected?",
    ("Result: “This is GATED — user must approve”", True),
])
add_text(s, W / 2 - 0.3, H * 0.535, 0.6, H * 0.03, u"↓", 16, NAVY, bold=True, align=PP_ALIGN.CENTER)
layer_box(0.57, 0.195, BG_YELLOWTEAL_LT, TEAL, u"✅", "Layer 3: User Controls", TEAL, [
    "Sees the proposal (the draft email)",
    "Approves or rejects (no auto-execute)",
    "4-minute undo window (can cancel if they change their mind)",
    "Full audit trail (can see why the code decided GATED)",
])
add_box(s, 1.0, H * 0.79, W - 2.0, H * 0.105, fill_hex=BG_GRAY_LT, line_hex="BDBDBD", line_w=1)
add_text(s, 1.2, H * 0.795, W - 2.4, H * 0.095,
         "User gets the help they need without the all-or-nothing risk. The AI proposes. Code decides. "
         "Trust is earned gradually. When unsure, it stops and asks. Reversible, not just approvals.",
         13, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
set_notes(s, "Here's how ARGUS works. Layer 1: GPT reads the user's command and proposes a structured "
             "action. Layer 2: Deterministic code decides if it's safe. Not the AI — the code. The "
             "policy engine checks trust, the hard-stop epoch validates freshness, the private-contact "
             "guard blocks sensitive addresses. Layer 3: User controls. They see the proposal, approve, "
             "and have a 4-minute window to undo. Every decision is audited. This separates intelligence "
             "(the LLM) from authority (the code). User gets scale without risk.")

# ════════════════════════════════ SLIDE 7 — ARCHITECTURE ════════════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.04, W, H * 0.08, "The Architecture: Why It's Fail-Safe",
         28, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)


def flow_box(x_pct, y_pct, w_pct, h_pct, bg, border, icon, title, title_color, sub=None):
    add_box(s, W * x_pct, H * y_pct, W * w_pct, H * h_pct, fill_hex=bg, line_hex=border, line_w=1.5)
    add_text(s, W * x_pct + 0.05, H * y_pct + 0.04, W * w_pct - 0.1, 0.22, icon + "  " + title,
              13, title_color, bold=True)
    if sub:
        add_bullets(s, W * x_pct + 0.08, H * y_pct + 0.30, W * w_pct - 0.16, H * h_pct - 0.34, sub,
                    10.5, GRAY12, line_spacing=1.15, space_after_px=1, bullet=False)


flow_box(0.05, 0.15, 0.16, 0.10, BG_BLUE_LT, BLUE, u"\U0001F4E7", "EMAIL → COMMAND", BLUE)
flow_box(0.25, 0.15, 0.20, 0.13, BG_BLUE_LT, BLUE, u"\U0001F9E0", "LAYER 1: GPT-4o", BLUE,
         ["(Interprets intent, generates proposal)"])
flow_box(0.18, 0.33, 0.64, 0.18, BG_GREEN_LT, GREEN, u"⚙", "LAYER 2: POLICY ENGINE", GREEN, [
    u"✓ Trust score check (action type + history)",
    u"✓ Hard-stop epoch validation (all approvals fresh?)",
    u"✓ Private-contact guard (is this address protected?)",
    u"→ Decision: ALLOW / GATED / BLOCK",
])
flow_box(0.18, 0.55, 0.64, 0.14, BG_YELLOWTEAL_LT, TEAL, u"✅", "LAYER 3: USER'S DECISION", TEAL, [
    "If ALLOW: Executes immediately with audit",
    "If GATED: Queue for approval + 4-min undo window",
])
flow_box(0.30, 0.73, 0.40, 0.10, BG_GRAY_LT, "BDBDBD", u"\U0001F4E7", "GMAIL API (Atomic, crash-safe)", DGRAY)


def badge(x_pct, icon, label, desc):
    add_text(s, W * x_pct, H * 0.875, W * 0.24, 0.05, icon + "  " + label, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, W * x_pct, H * 0.925, W * 0.24, 0.07, desc, 10, GRAY12, align=PP_ALIGN.CENTER, line_spacing=1.1)


badge(0.05, u"\U0001F510", "Atomic", "All-or-nothing transactions (no partial states)")
badge(0.38, u"\U0001F504", "Reversible", "4-minute undo window for all actions")
badge(0.71, u"\U0001F4DD", "Audited", "Every decision logged with cryptographic proof")
set_notes(s, "The architecture has three layers. Layer 1 is GPT — it interprets what the user wants. "
             "Layer 2 is where the magic happens: deterministic code decides if it's safe. Not a guess, "
             "not a probability — a decision. It checks trust, validates freshness, blocks protected "
             "addresses. Layer 3 is the user. If the code says ALLOW, it goes. If GATED, they approve. "
             "They always have a 4-minute window to undo. And everything is cryptographically audited "
             "so they can prove what happened.")

# ════════════════════════════════ SLIDE 8 — LIVE DEMO ═══════════════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.04, W, H * 0.08, "See It In Action", 30, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)
captions = [
    "1. Inbox — select an email from a colleague",
    "2. Command — type your intent",
    "3. Proposal — ARGUS generates draft",
    "4. Queue — approve + 4:00 countdown",
    "5. Executed — toast confirmation",
]
slot_w = (W - 1.2) / 5
for i, cap in enumerate(captions):
    x = 0.6 + i * slot_w
    add_box(s, x + 0.08, H * 0.22, slot_w - 0.16, H * 0.40, fill_hex=BG_GRAY_LT, line_hex="CCCCCC", line_w=1)
    add_text(s, x + 0.08, H * 0.40, slot_w - 0.16, 0.3, "[screenshot]", 11, GRAY12, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.05, H * 0.63, slot_w - 0.10, H * 0.10, cap, 12, NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 4:
        add_text(s, x + slot_w - 0.12, H * 0.38, 0.24, 0.10, u"→", 16, NAVY, bold=True, align=PP_ALIGN.CENTER)
add_box(s, W * 0.20, H * 0.76, W * 0.60, H * 0.13, fill_hex=BG_GREEN_LT, line_hex=GREEN, line_w=1)
add_text(s, W * 0.22, H * 0.775, W * 0.56, H * 0.10,
         "Live Demo URL: localhost:8081  ·  Credentials: PROJECT_ARGUS / ARGUS_DEMO\n"
         "Features Available: Emergency stop, profile switching, audit trail inspection",
         12, DGRAY, font=MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
add_text(s, 0, H * 0.93, W, H * 0.06,
         "Live demo is running during the pitch. If connection fails, screenshots above are your fallback.",
         11, GRAY12, italic=True, align=PP_ALIGN.CENTER)
set_notes(s, "Let me show you ARGUS working. The user selects an email, types a command, ARGUS proposes "
             "a draft, they see it's GATED and requires approval, they approve, and they get a 4-minute "
             "window to undo if they change their mind. Then it sends. Full transparency, full control, "
             "zero surprises.\n\n[BUILDER NOTE: drop real screenshots from the running demo into the 5 "
             "placeholder boxes above before presenting — none were available to embed automatically.]")

# ════════════════════════════════ SLIDE 9 — VALIDATION (SURVEY) ═════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.05, W, H * 0.08, "What Users Actually Want (Survey: N=17)",
         28, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)


def stat_card(x_pct, bg, border, num, num_color, desc, sub):
    add_box(s, W * x_pct, H * 0.22, W * 0.28, H * 0.42, fill_hex=bg, line_hex=border, line_w=1.75)
    add_text(s, W * x_pct + 0.15, H * 0.255, W * 0.28 - 0.30, 0.05, u"✅", 16, TEAL, align=PP_ALIGN.RIGHT)
    add_text(s, W * x_pct, H * 0.30, W * 0.28, H * 0.14, num, 40, num_color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, W * x_pct + 0.15, H * 0.45, W * 0.28 - 0.30, H * 0.10, desc, 13, DGRAY, align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_text(s, W * x_pct + 0.15, H * 0.555, W * 0.28 - 0.30, H * 0.07, sub, 11, GRAY12, italic=True, align=PP_ALIGN.CENTER)


stat_card(0.06, BG_ORANGE_LT, BORDER_ORANGE, "88%", AMBER, "prefer “ask permission every time”", "vs “earn autonomy gradually”")
stat_card(0.36, BG_GREEN_LT, BORDER_GREEN, "100%", GREEN, "want AI to pause when unsure", "vs make its best guess and act")
stat_card(0.66, BG_BLUE_LT, BORDER_BLUE, "100%", BLUE, "want to see exactly why each decision was made", "(survey quote: “justify its every decision”)")
add_box(s, W * 0.10, H * 0.74, W * 0.80, H * 0.16, fill_hex=BG_GRAY_LT, line_hex=TEAL, line_w=1.75)
add_text(s, W * 0.13, H * 0.755, W * 0.74, H * 0.13,
         "Users didn't ask for a ‘nice AI email tool.’ They asked for exactly what ARGUS builds. "
         "We didn't invent demand — we solved for it.",
         15, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
set_notes(s, "We surveyed 17 users about AI email management. 88% want permission-based control, not "
             "gradual autonomy. 100% want AI to pause when unsure. 100% want to see why every decision "
             "was made. Those aren't nice-to-haves. They're non-negotiable. Users didn't ask for a fancy "
             "AI email tool. They asked for ARGUS. We built what they're asking for.")

# ════════════════════════════════ SLIDE 10 — PROOF & CLOSE ══════════════════
s = new_slide(WHITE)
add_text(s, 0, H * 0.05, W, H * 0.08, "We Shipped Real Code", 30, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER)

stats = [
    ("869/869", TEAL, "Tests Passing (100% coverage)"),
    ("7", GREEN, "Fail-Safe Controls\n(hard-stop, rate limit, dedup, epoch,\nprivate contacts, audit chain, stale detection)"),
    ("3", BLUE, "Angles of Testing\n(Normal flow + Hacker/adversarial + Strict validation)"),
    ("0", RED, "Silent Failures\n(cryptographic audit guarantees transparency)"),
]
y = 0.24
for num, color, label in stats:
    add_text(s, 0.6, H * y, 1.7, 0.30, num, 34, color, bold=True)
    add_text(s, 2.4, H * y + 0.01, 3.9, 0.30, label, 11.5, DGRAY, line_spacing=1.15, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.155

add_box(s, W * 0.55, H * 0.22, W * 0.40, H * 0.40, fill_hex=BG_GREEN_LT, line_hex=GREEN, line_w=1.75)
add_text(s, W * 0.58, H * 0.255, W * 0.34, 0.05, "Ready To Ship", 16, GREEN, bold=True)
tb = s.shapes.add_textbox(Inches(W * 0.58), Inches(H * 0.33), Inches(W * 0.34), Inches(H * 0.27))
tf = tb.text_frame
tf.word_wrap = True
for i, item in enumerate(["Policy engine locked (Phase 8 complete)", "Live demo running",
                          "Full test suite passing", "Architecture scales beyond email"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = u"✅ " + item
    r.font.size = Pt(13)
    r.font.name = SANS
    r.font.color.rgb = rgb(DGRAY)

add_text(s, W * 0.10, H * 0.68, W * 0.80, H * 0.13,
         "AI should amplify your judgment, not replace it. ARGUS is that AI.",
         22, NAVY, font=SERIF, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

links = [
    (u"\U0001F517 Live Demo: localhost:8081"),
    (u"\U0001F4BE GitHub: [repo link]"),
    (u"\U0001F4E7 Contact: kayden.low24@gmail.com"),
]
lx = 0.6
lw = (W - 1.2 - 0.4) / 3
for i, link in enumerate(links):
    bx = lx + i * (lw + 0.2)
    add_box(s, bx, H * 0.86, lw, H * 0.09, fill_hex=BG_TEAL_LT, line_hex=TEAL, line_w=1)
    add_text(s, bx, H * 0.865, lw, H * 0.08, link, 11.5, TEAL, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "We didn't build a prototype. ARGUS has 869 passing tests, three angles of validation, "
             "seven fail-safe controls, and zero silent failures. The policy engine is locked. The demo "
             "is live. The architecture is ready to scale. We've solved the bottleneck: LLMs propose, "
             "code decides, trust is earned. This isn't a nice idea. It's production-ready software "
             "that solves the exact problem judges are asking about: how do you trust an AI to act? "
             "You make the code decide, not the AI. That's ARGUS.\n\n[BUILDER NOTE: replace '[repo link]' "
             "with the real GitHub URL before submission.]")

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
