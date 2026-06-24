"""Slide 10 currently repeats the Level-4 ICP that slide 3 already owns
(target user & use case). Fix: drop the persona recap from slide 10's left
column and replace it with content that's actually new — what 'real product'
means beyond the hackathon demo (platform/market expansion, not user
identity). Right column (Today / architecture / rebuild / ask / why now) and
the bottom band already avoid repetition, so they're untouched. Title is
re-trimmed to drop its persona callback too. Design (position/size/font/
color) is untouched — text only.
"""
from pptx import Presentation

SRC = r"C:\Users\baldwin\Downloads\dck.pptx"
OUT = r"C:\Users\baldwin\Downloads\dck_v2.pptx"

prs = Presentation(SRC)
s = prs.slides[9]


def shape(sid):
    for sh in s.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_runs(sid, texts):
    tf = shape(sid).text_frame
    p0 = tf.paragraphs[0]
    for r, t in zip(p0.runs, texts):
        r.text = t
    for r in list(p0.runs[len(texts):]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def set_text(sid, text):
    set_runs(sid, [text])


# Title: drop the persona callback, keep it pure vision/ask
set_runs(273, [
    "The Vision Beyond the Hackathon.",
    " Here's the Ask.",
])

# Left column: replace the repeated ICP with new content — what "real product" means
set_text(274, "WHAT 'REAL PRODUCT' MEANS")
set_text(275,
    "A hackathon weekend proved the core mechanic works. A real product means "
    "every inbox provider, not just Gmail — Outlook, IMAP, any account an agent "
    "can act through. It means every tool an agent touches, not just email — "
    "calendars, CRMs, file systems, anything with an irreversible action. One "
    "deterministic permission and trust layer underneath all of it, so every "
    "team shipping agents stops rebuilding the same safety net from scratch. "
    "That is the company-shaped version of what we built this weekend.")

# Right column (278-290) and bottom band (292/293) already avoid repeating
# slide 3's ICP — left untouched.

prs.save(OUT)
print("Saved:", OUT)
