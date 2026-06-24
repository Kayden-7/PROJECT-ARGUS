"""Apply the same slide-10 content swap (Vision & The Ask) to the user's own
redesigned copy of the deck. This file's text shapes are single-color (no
alternating emphasis runs like the earlier deck) — content only, no design
touched: no color/font/size/position changes, just the words.
"""
from pptx import Presentation

SRC = r"C:\Users\baldwin\Downloads\ARGUS_pitch_deck(1)(1)(1).pptx"
OUT = r"C:\Users\baldwin\Downloads\ARGUS_pitch_deck(1)(1)(1)_v2.pptx"

prs = Presentation(SRC)
s = prs.slides[9]


def shape(sid):
    for sh in s.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_runs(sid, texts):
    tf = shape(sid).text_frame
    p0 = tf.paragraphs[0]
    assert len(p0.runs) >= len(texts), f"{sid}: only {len(p0.runs)} runs for {len(texts)} texts"
    for r, t in zip(p0.runs, texts):
        r.text = t
    for r in list(p0.runs[len(texts):]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def set_text(sid, text):
    set_runs(sid, [text])


# Title (2 runs)
set_runs(273, [
    "Built for the One Email You Can't Get Wrong.",
    " Now, the Vision Beyond It.",
])

# Left column
set_text(274, "WHO WE BUILT THIS FOR")
set_text(275,
    "A solo founder or freelance consultant whose livelihood runs on email — "
    "who wants an AI assistant to handle the inbox, but won't hand it over "
    "because one wrong auto-sent message to a client, or one deleted thread, "
    "could cost a relationship or a contract. So today, they still do it all "
    "by hand — and ARGUS exists so they finally don't have to.")

# Right column
set_text(276, "WHERE ARGUS GOES FROM HERE")

set_text(278, "Today")
set_text(279, "A real, working product, not a slide — live Gmail OAuth, a deterministic policy engine, and 863 passing automated tests.")

set_text(281, "The architecture already plans for this")
set_text(282, "The PROPOSE layer was built model-agnostic — swap the model, the deterministic core stays the same.")

set_text(284, "The rebuild")
set_text(285, "The next build is Claude-native — Claude reasons, the deterministic core still decides every action, unchanged.")

set_text(287, "The ask")
set_text(288, "Help us take ARGUS past the hackathon: harden the Claude-native build for real inboxes.")

set_text(289, "Why now")
set_text(290, "Every agent platform hits the same wall: authority without losing control. We intend to build the layer that solves it.")

# Bottom band
set_text(292, "This began as a hackathon build. The vision: a real product, Claude-native, trusted by the people who need it.")
# 293 ("AI proposes. Code decides.") stays unchanged.

prs.save(OUT)
print("Saved:", OUT)
