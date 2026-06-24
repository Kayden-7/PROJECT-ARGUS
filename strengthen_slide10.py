"""Strengthen slide 10's wording — sharper hooks, more concrete claims, a
harder ask, punchier closes. Builds on dck_v2.pptx (which already fixed the
slide-3 ICP repetition). Text only — no shape/position/font/color changes.
"""
from pptx import Presentation

SRC = r"C:\Users\baldwin\Downloads\dck_v2.pptx"
OUT = r"C:\Users\baldwin\Downloads\dck_v3.pptx"

prs = Presentation(SRC)
s = prs.slides[9]


def shape(sid):
    for sh in s.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_runs(sid, texts):
    tf = shape(sid).text_frame
    p0 = tf.paragraphs[0]
    for r, t in zip(p0.runs, texts):
        r.text = t
    for r in list(p0.runs[len(texts):]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def set_text(sid, text):
    set_runs(sid, [text])


# Title
set_runs(273, [
    "This Weekend Proved the Mechanic Works.",
    " Now, the Real Ask.",
])

# Left column — what "real product" means, sharpened
set_text(274, "FROM WEEKEND BUILD TO INDUSTRY LAYER")
set_text(275,
    "A hackathon weekend proved the mechanic works. A real product means "
    "every inbox provider, not just Gmail — Outlook, IMAP, any account an agent "
    "can act through. It means every irreversible action an agent takes, not "
    "just email — calendars, CRMs, file systems, payments. One deterministic "
    "permission and trust layer underneath all of it, so every team shipping "
    "agents stops rebuilding the same safety net from scratch, badly, alone. "
    "That's not a feature. That's infrastructure.")

# Right column — sharper claims, harder ask, bigger insight
set_text(279, "Not a slide, not a mockup — live Gmail OAuth, a deterministic policy engine, and 863 automated tests, all green.")
set_text(282, "We built the PROPOSE layer swappable on day one — change the model, the core that decides never moves.")
set_text(285, "Next: Claude-native. Claude does the reasoning; the same deterministic core still has final say — intelligence was never the part we needed to fix.")
set_text(288, "Back us to take ARGUS past the hackathon — fund the Claude-native build and put it in front of real inboxes, not just judges.")
set_text(290, "Every agent platform will hit this same wall: real authority without losing control. Whoever builds that layer first becomes infrastructure. We intend to be first.")

# Bottom band — bigger close
set_text(292, "This began as a hackathon build. The vision is bigger: the trust layer every autonomous agent runs through — Claude-native, real, and shipped.")
# 293 ("AI proposes. Code decides.") stays unchanged.

prs.save(OUT)
print("Saved:", OUT)
