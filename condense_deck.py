"""Condense the designed 12-slide deck to 10 by merging the two natural pairs,
preserving every slide's theme/design. No content is dropped — only merged.

  Pair A: Solution & Value (S7, dark) + Why ARGUS Wins (S10, light)
          -> kept on S7's dark design; its right column (3 guarantees, a subset
             of Why-Wins) is replaced by the full 6 defensibility points.
  Pair B: How It Works architecture (S8) + three outcomes (S9), both light
          -> kept on S8; the ALLOW/GATED/BLOCK examples + earned-trust fold into
             S8's existing DECIDE / EXECUTE boxes.

Then S9 and S10 are removed (their content now lives on S8 and S7), and page /
section numbers are renumbered. Other slides are byte-identical.
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC = r"C:\Users\baldwin\Downloads\Website_redesign_deck(2).pptx"
OUT = r"C:\Users\baldwin\Downloads\Website_redesign_deck_10slides.pptx"
prs = Presentation(SRC)
slides = list(prs.slides)

# theme colors (extracted from the deck)
CREAM = RGBColor.from_string("EDEBE4")
LIGHTGREY = RGBColor.from_string("D7DAE1")
GREEN_D = RGBColor.from_string("7FB79B")   # green for dark slides
LABEL = RGBColor.from_string("A9B2C2")


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_text(slide, name, new_text):
    sh = shape_by_name(slide, name)
    if not sh or not sh.has_text_frame:
        return False
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = new_text
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    return True


def remove_shapes(slide, names):
    for sh in list(slide.shapes):
        if sh.name in names:
            sh._element.getparent().remove(sh._element)


def add_text(slide, x, y, w, h, text, size, color, bold=False, font="Arial"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

# ════════════════════ PAIR A — merge Why-Wins (S10) into Solution (S7) ════════════════════
s7 = slides[6]
set_text(s7, "Text 6", "Solution, Value & Why It Wins")
# remove the 3-guarantee block (right column) — it's a subset of the 6 points
remove_shapes(s7, ["Text 13", "Shape 14", "Text 15", "Text 16",
                   "Shape 17", "Text 18", "Text 19", "Shape 20", "Text 21", "Text 22"])
# add the full 6 defensibility points in the same right column, dark theme
add_text(s7, 10.7, 3.55, 8.5, 0.22, "WHY IT'S DEFENSIBLE", 10.5, LABEL, bold=False, font="Courier New")
points = [
    ("PROMPT INJECTION CAN'T MOVE THE DECISION",
     "Code re-derives permission independently — fully compromise the model and it still cannot send an unapproved external email."),
    ("DEFENCE IN DEPTH AGAINST INJECTION",
     "Body-only drafting context (no recipient authority), allow-listed style with no free-form field, and every external send gated."),
    ("FAILS CLOSED BY CONSTRUCTION",
     "Any uncertainty routes to MANUAL_REVIEW instead of a guess — safe is the default path, not the lucky one."),
    ("EARNED TRUST, NOT DECLARED TRUST",
     "Autonomy grows with proven reliability per action type and collapses fast after a failure. No “just trust me.”"),
    ("PROVABLE, NOT MERELY SAFE",
     "Append-only, SHA-256 hash-chained audit log with a verify endpoint and full replay — tamper-evident by design."),
    ("MODEL-AGNOSTIC INFRASTRUCTURE",
     "Swap GPT-4o for any model and the deterministic control plane is unchanged — a reusable layer, not one app."),
]
y = 3.95
for hdr, body in points:
    add_text(s7, 10.7, y, 8.5, 0.22, hdr, 11.0, GREEN_D, bold=True, font="Courier New")
    add_text(s7, 10.7, y + 0.27, 8.5, 0.6, body, 11.0, LIGHTGREY, bold=False, font="Arial")
    y += 0.95

# ════════════════════ PAIR B — merge three outcomes (S9) into How It Works (S8) ════════════════════
s8 = slides[7]
set_text(s8, "Text 6", "How It Works— architecture & the three outcomes it produces")
set_text(s8, "Text 7",
         "Three strictly separated layers, and the same pipeline producing three outcomes — every decision made by code, never the model.")
# DECIDE box: keep the gauntlet, append the worked ALLOW/GATED/BLOCK examples
set_text(s8, "Text 17",
         "Every proposal runs the same gauntlet, every time: (a) prime-rule check — hard BLOCKs first; "
         "(b) action taxonomy — 9 “free” actions auto-allowed, 11 “gated” actions always need a human; "
         "(c) earned-trust check vs the active profile threshold (Strict / Balanced / Autonomous); "
         "(d) safety filter — can only downgrade Allow to Gate (every external send is gated regardless of trust). "
         "Output: ALLOW / GATED / BLOCK, always with a reason trace. "
         "Worked examples — ALLOW: “mark the newsletter as read” (free, instant); "
         "GATED: “reply to my client” (external send → human approval with a live countdown); "
         "BLOCK: “delete every email from my boss” (prime rule → hard stop, nothing happens).")
# EXECUTE box: keep crash-safe note, append the earned-trust mechanics
set_text(s8, "Text 22",
         "On ANY uncertainty — a crash mid-send, an ambiguous state — it fails closed to MANUAL_REVIEW; never a double-send, never a lost email. "
         "Earned trust: starts at 40/100 and rises only after a verified send (approval alone never counts); after a high-severity failure, gains halve for a 10-event window — autonomy earned slowly, revoked fast.")

# ════════════════════ remove S9 (outcomes) and S10 (why-wins) ════════════════════
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for idx in sorted([8, 9], reverse=True):  # 0-based: S9, S10
    sld = ids[idx]
    rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sld)

# ════════════════════ renumber page indicators + section numbers ════════════════════
page_re = re.compile(r"^\d{1,2}\s*/\s*\d{1,2}$")
sect_re = re.compile(r"^\d{2}$")
total = len(list(prs.slides))
for ordinal, s in enumerate(prs.slides, start=1):
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if page_re.match(t) and "/" in t:
            set_text(s, sh.name, f"{ordinal:02d} / {total}")
        elif sect_re.match(t):
            set_text(s, sh.name, f"{ordinal-1:02d}")

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
